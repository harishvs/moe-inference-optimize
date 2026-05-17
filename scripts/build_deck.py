"""Generate `talk/deck.pptx` for the AI Infrastructure presentation.

Audience: Technical Domain Expert + VP of Strategy. Format: 30-minute slot,
~15 minutes of slides + ~15 minutes of Q&A. Numbers are pulled live from
`results/` via `talk.numbers`.

Layout: clean modern minimal (16:9), no template inheritance — every shape is
positioned explicitly so there's no fighting with master-slide artifacts.

Run:
    uv run python scripts/build_deck.py
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

REPO = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO))
from talk import numbers  # noqa: E402

FIG_DIR = REPO / "talk" / "figures"
OUT = REPO / "talk" / "deck.pptx"

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Restrained palette
INK = RGBColor(0x1F, 0x2A, 0x37)        # near-black title
SUBTLE = RGBColor(0x55, 0x5F, 0x6E)     # body grey
ACCENT = RGBColor(0x2E, 0x66, 0xC2)     # cool blue
WARN = RGBColor(0xC0, 0x4A, 0x2A)       # rust for caveats
RULE = RGBColor(0xE2, 0xE5, 0xEA)       # divider grey
BG = RGBColor(0xFF, 0xFF, 0xFF)         # white
TABLE_HEAD_BG = RGBColor(0xF2, 0xF4, 0xF7)
TABLE_ALT_BG = RGBColor(0xFA, 0xFB, 0xFC)

FONT = "Calibri"


# ------------------------------------------------------------
# Primitives
# ------------------------------------------------------------

def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # "blank"


def add_text(slide, text: str, *, left: float, top: float, width: float, height: float,
             pt: int = 18, bold: bool = False, color: RGBColor = INK,
             align: PP_ALIGN = PP_ALIGN.LEFT, italic: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    for r in p.runs:
        r.font.name = FONT
        r.font.size = Pt(pt)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color


def add_rule(slide, *, left: float, top: float, width: float, color: RGBColor = ACCENT,
             height_pt: float = 3.0) -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(left), Inches(top),
                                  Inches(width), Pt(height_pt))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def add_title(slide, title: str, *, eyebrow: str | None = None) -> None:
    if eyebrow:
        add_text(slide, eyebrow, left=0.6, top=0.45, width=12.2, height=0.3,
                 pt=12, bold=True, color=ACCENT)
        add_text(slide, title, left=0.6, top=0.78, width=12.2, height=0.7,
                 pt=30, bold=True, color=INK)
    else:
        add_text(slide, title, left=0.6, top=0.5, width=12.2, height=0.7,
                 pt=30, bold=True, color=INK)
    add_rule(slide, left=0.6, top=1.55, width=1.0)


def add_footer(slide, page_num: int) -> None:
    add_text(slide, f"OLMoE-1B-7B (1B active / 7B total) on Blackwell  ·  diagnosis & optimization",
             left=0.6, top=7.05, width=10.0, height=0.3,
             pt=10, color=SUBTLE)
    add_text(slide, f"{page_num}",
             left=12.6, top=7.05, width=0.6, height=0.3,
             pt=10, color=SUBTLE, align=PP_ALIGN.RIGHT)


def add_bullets(slide, bullets: list[str], *, left: float = 0.6, top: float = 1.9,
                width: float = 12.2, height: float = 4.5,
                pt: int = 20, line_gap_pt: int = 8) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + b
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(line_gap_pt)
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(pt)
            r.font.color.rgb = INK


def add_table(slide, rows: list[list[str]], *, left: float, top: float,
              width: float, height: float, pt: int = 14,
              col_widths: list[float] | None = None,
              cell_text_colors: dict[tuple[int, int], RGBColor] | None = None
              ) -> None:
    n_rows, n_cols = len(rows), len(rows[0])
    shape = slide.shapes.add_table(n_rows, n_cols,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tbl = shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)
    cell_text_colors = cell_text_colors or {}
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            cell.margin_left = Emu(60000)
            cell.margin_right = Emu(60000)
            cell.margin_top = Emu(40000)
            cell.margin_bottom = Emu(40000)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = TABLE_HEAD_BG
            elif r % 2 == 0:
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            else:
                cell.fill.fore_color.rgb = BG
            override = cell_text_colors.get((r, c))
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = FONT
                    run.font.size = Pt(pt)
                    run.font.bold = (r == 0) or (override is not None)
                    if override is not None:
                        run.font.color.rgb = override
                    else:
                        run.font.color.rgb = INK if r > 0 else SUBTLE


def add_caption(slide, text: str, *, top: float, color: RGBColor = SUBTLE,
                pt: int = 14) -> None:
    add_text(slide, text, left=0.6, top=top, width=12.2, height=0.6,
             pt=pt, italic=True, color=color)


def add_picture(slide, image_path: Path, *, left: float, top: float,
                width: float, height: float | None = None) -> None:
    if not image_path.exists():
        add_text(slide, f"[figure missing: {image_path.name}]",
                 left=left, top=top, width=width, height=0.6,
                 pt=14, color=WARN, italic=True)
        return
    kwargs = {"width": Inches(width)}
    if height is not None:
        kwargs["height"] = Inches(height)
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), **kwargs)


def set_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def fmt_pct(v: float | None, digits: int = 0) -> str:
    return f"{v:+.{digits}f}%" if v is not None else "TBD"


def fmt_ms(v: float | None, digits: int = 1) -> str:
    return f"{v:.{digits}f}" if v is not None else "TBD"


# ------------------------------------------------------------
# Slides
# ------------------------------------------------------------

def slide_title(prs: Presentation) -> None:
    s = blank(prs)
    add_text(s, "AI Infrastructure",
             left=0.8, top=2.1, width=12.0, height=0.5,
             pt=14, bold=True, color=ACCENT)
    add_text(s, "Ship the chat assistant",
             left=0.8, top=2.5, width=12.0, height=1.0,
             pt=48, bold=True, color=INK)
    add_text(s, "Cutting per-request response time on OLMoE-1B-7B / vLLM",
             left=0.8, top=3.7, width=12.0, height=0.5,
             pt=20, color=SUBTLE)
    add_rule(s, left=0.8, top=4.7, width=2.0, height_pt=4.0)
    add_text(s, "Harish Rao",
             left=0.8, top=5.0, width=12.0, height=0.5,
             pt=20, bold=True, color=INK)
    add_text(s, "Cloud AI Infrastructure Engineer",
             left=0.8, top=5.5, width=12.0, height=0.4,
             pt=14, color=SUBTLE)
    set_notes(s,
        "Open by anchoring the audience in the actual product situation: a "
        "chat application running on EKS that's missing its inference latency "
        "latency target under load. The model is OLMoE-1B-7B-Instruct served by vLLM. "
        "We're going to diagnose it the same way we'd diagnose a production "
        "incident — measure, profile, list the levers, attribute the deltas, "
        "and ship a ranked plan with measured tradeoffs.\n\n"
        "Two audiences: the technical expert who'll probe kernel-level "
        "tradeoffs, and the VP who needs to know what this means for cost, "
        "time-to-market, and the 10× scenario. Detail lives in the appendix "
        "and in the repo."
    )


def slide_problem(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "The problem in one slide")

    # Left: deployment topology diagram
    add_text(s, "DEPLOYMENT", left=0.6, top=1.85, width=7.2, height=0.3,
             pt=11, bold=True, color=ACCENT)
    add_picture(s, FIG_DIR / "diagram_eks_topology.png",
                left=0.5, top=2.2, width=7.4)

    # Right: what's broken
    add_text(s, "WHAT'S BROKEN", left=8.3, top=1.85, width=4.6, height=0.3,
             pt=11, bold=True, color=ACCENT)
    add_bullets(s, [
        "Baseline misses the per-request SLA from prompts ≥ 512 tokens — "
        "exactly the chat-app range",
        "User pain: noticeable lag once prompts include the system prompt + "
        "tool descriptions + history",
        "Business pain: launch blocked on a sub-second instant-feel experience — "
        "adding GPUs doesn't help, latency is per-request",
        "Question: which levers cut per-request latency, by how much, and at "
        "what cost?",
    ], left=8.3, top=2.25, width=4.6, height=4.5, pt=13, line_gap_pt=6)
    add_footer(s, 2)
    set_notes(s,
        "Frame the problem so both panelists see themselves in it. The "
        "technical expert hears 'inference latency, MoE model, vLLM, EKS' "
        "and thinks kernels and routing. The VP hears 'missing the latency target, blocks "
        "launch, over-provisioning is a tax' and thinks cost and "
        "time-to-market.\n\nDon't dwell — the rest of the talk is the answer."
    )


def slide_system(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "What 'fast' means for this app", eyebrow="DEFINITIONS")

    # Left: timeline diagram
    add_text(s, "ONE REQUEST, OVER TIME", left=0.6, top=1.85, width=7.2, height=0.3,
             pt=11, bold=True, color=ACCENT)
    add_picture(s, FIG_DIR / "diagram_ttft_tpot_timeline.png",
                left=0.4, top=2.3, width=7.6)

    # Right: what we measure
    add_text(s, "WHAT I MEASURED", left=8.3, top=1.85, width=4.6, height=0.3,
             pt=11, bold=True, color=ACCENT)
    add_bullets(s, [
        "TTFT — time to first token (how long until anything appears)",
        "TPOT — time per output token (how fast the rest stream in)",
        "Total response time = TTFT + tokens × TPOT",
        "Test grid: typical chat-app prompt lengths (128–1500 tokens)",
        "Model max context: 4,096 tokens — test sizes stay well inside that",
        "SLA: TTFT < 15 ms, TPOT < 2.5 ms, end-to-end < 500 ms (p90, single-user)",
    ], left=8.3, top=2.25, width=4.6, height=4.3, pt=13, line_gap_pt=6)
    add_text(s,
             "SLA chosen for an \"instant-feel\" chat experience: TTFT below "
             "perceptual-immediacy threshold, streaming faster than reading "
             "speed, half-second total response. Pending product confirmation.",
             left=8.3, top=6.45, width=4.6, height=0.6,
             pt=9, italic=True, color=SUBTLE)
    add_footer(s, 3)
    set_notes(s,
        "Before measuring, the audience needs the same vocabulary we have. "
        "TTFT and TPOT are the two metrics any inference latency target is built on. "
        "TTFT is what the user feels right after they hit enter. TPOT is "
        "what they feel as the response streams in.\n\n"
        "The grid we measure on is deliberately chat-app-shaped: not "
        "synthetic, not maximum-context, but realistic for a customer-facing "
        "assistant with a system prompt and tool descriptions in the few-"
        "hundred-token range."
    )


def slide_method(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "The method", eyebrow="APPROACH")
    add_bullets(s, [
        "1.  Measure — clean, repeatable baseline (TTFT/TPOT, seed-pinned)",
        "2.  Profile — kernel-level breakdown via torch.profiler",
        "3.  Try one lever at a time — measure each delta against the baseline",
        "4.  Decide — rank by impact, name the trade-offs, recommend",
    ])
    add_caption(s, "Skip any step and I'd optimize the wrong thing — or claim a win that's just noise.",
                top=6.3, color=WARN, pt=15)
    add_footer(s, 4)
    set_notes(s,
        "This is the meta-point that holds the whole talk together. Most "
        "performance work goes wrong because someone skips one of these "
        "steps. They jump to 'I'll write a kernel' before profiling, or "
        "they 'optimize' something and forget to measure variance, or they "
        "ship a 5× win that's actually noise from a different prompt.\n\n"
        "Every claim in the rest of this talk follows this loop."
    )


def slide_baseline(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "Where I'm starting from — single Blackwell GPU, batch size = 1",
              eyebrow="MEASURE")
    bf = numbers.derive_tpot(numbers.latency_grid("baseline"))
    SLA_TTFT_MS = 15.0   # p90
    SLA_E2E_MS = 500.0   # 200-token response
    SLA_TPOT_MS = 2.5    # per token
    cell_colors: dict[tuple[int, int], RGBColor] = {}
    if bf:
        rows = [["Input length", "TTFT mean (ms)", "TTFT p90 (ms)",
                 "End-to-end @ 200 tokens (ms)", "TPOT (ms per token)"]]
        for r_idx, (L, v) in enumerate(sorted(bf.items()), start=1):
            rows.append([
                str(L),
                fmt_ms(v["ttft_mean_ms"]),
                fmt_ms(v["ttft_p90_ms"]),
                fmt_ms(v["end_to_end_mean_ms"], 0),
                f"{v['tpot_mean_ms']:.2f}",
            ])
            # Flag cells that miss the SLA
            if v["ttft_mean_ms"] >= SLA_TTFT_MS:
                cell_colors[(r_idx, 1)] = WARN
            if v["ttft_p90_ms"] >= SLA_TTFT_MS:
                cell_colors[(r_idx, 2)] = WARN
            if v["end_to_end_mean_ms"] >= SLA_E2E_MS:
                cell_colors[(r_idx, 3)] = WARN
            if v["tpot_mean_ms"] >= SLA_TPOT_MS:
                cell_colors[(r_idx, 4)] = WARN
        add_table(s, rows, left=0.6, top=2.0, width=12.1, height=2.8, pt=15,
                  col_widths=[1.8, 2.4, 2.3, 3.2, 2.4],
                  cell_text_colors=cell_colors)
    add_caption(s,
        "Red cells miss the SLA (TTFT < 15 ms, TPOT < 2.5 ms, e2e < 500 ms). "
        "Baseline starts failing at L=512 — exactly the chat-app range.",
        top=5.1)
    add_footer(s, 5)
    set_notes(s,
        "Three things the audience should notice:\n\n"
        "First: variance is extremely tight across 30 iterations. p90 sits "
        "less than a millisecond from the mean. This is the reward for "
        "fixing the seed and running enough iterations — any speedup we "
        "measure later is real, not an artifact of changing prompts or "
        "warm caches.\n\n"
        "Second: TTFT grows sub-linearly with input length, not "
        "quadratically. The N² attention term is small at this scale; "
        "FFN dominates.\n\n"
        "Third: TPOT is nearly flat across input length. Decode doesn't "
        "care how long the prompt was. Prefill is compute-bound; decode is "
        "bandwidth-bound. Two different problems."
    )


def slide_profile(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "Where the time goes — and where it doesn't", eyebrow="PROFILE")

    # Left: kernel breakdown bar
    add_text(s, "PREFILL TTFT BREAKDOWN", left=0.6, top=1.85, width=7.2, height=0.3,
             pt=11, bold=True, color=ACCENT)
    add_picture(s, FIG_DIR / "diagram_prefill_kernels.png",
                left=0.5, top=2.4, width=7.6)

    # Right: takeaways
    add_text(s, "TAKEAWAYS", left=8.3, top=1.85, width=4.6, height=0.3,
             pt=11, bold=True, color=ACCENT)
    add_bullets(s, [
        "fused_moe_kernel — the expert FFN matmuls, one fused kernel per layer (tensor cores) — owns ~68% of prefill",
        "Attention projection GEMMs are second at ~18%",
        "Hot path is already production-tuned in vLLM — no kernel headroom on the dominant paths",
        "Levers live above the kernels: quantization, caching, batching, parallelism",
    ], left=8.3, top=2.25, width=4.6, height=4.5, pt=14, line_gap_pt=6)
    add_footer(s, 6)
    set_notes(s,
        "Walk the audience across the bar. Numbers are from torch.profiler "
        "on Blackwell at L=1024, bf16, single-user (results/profile/"
        "baseline_20260516_000526_L1024_seed42/), renormalized to 100% of "
        "accounted leaf-kernel GPU time.\n\n"
        "fused_moe_kernel — the grouped-GEMM kernel that runs the MoE "
        "feed-forward — owns 68% of prefill. It's already tuned for tensor "
        "cores; vLLM uses flashinfer's fused-MoE path on Blackwell. "
        "Attention projection GEMMs (CUTLASS bf16) are second at 18%, "
        "with attention math itself (flash_fwd_splitkv) at ~5% and MoE "
        "expert dispatch/combine kernels at another ~5%. The router "
        "(top-K softmax) is under 0.5% — basically noise.\n\n"
        "Compared to L4: the relative shape is similar but Blackwell's "
        "tensor cores chew through the MoE FFN faster, so its share "
        "drops from L4's ~80% to ~68%, and attention projections rise "
        "into the visible range. Same conclusion either way: the hot "
        "path is FFN, FFN is already tuned, no kernel headroom for a "
        "solo developer.\n\n"
        "The point: the profile tells the engineering budget where to go. "
        "It pointed at FP8 (attacks fused_moe via lower bytes), prefix "
        "caching (skips prefill entirely), batching (amortizes weight "
        "loads), and TP=2 (splits the work across both GPUs). Those are "
        "the levers on the next slide.\n\n"
        "DECODE PROFILE (if asked) — different shape, same dominant "
        "kernel. From the decode trace at I=1, O=128 "
        "(results/profile/baseline_decode_20260516_000621_I1_O128_seed42/):\n\n"
        "  fused_moe_kernel:               33%\n"
        "  CUTLASS bf16 WMMA GEMMs:        14%   (attention QKV/O)\n"
        "  flash_fwd_splitkv:               4%   (attention math)\n"
        "  MoE dispatch/combine:            4%\n"
        "  MoE router (top-K):              2%\n"
        "  RMS norms:                       3%\n"
        "  KV cache writes:                 1%\n"
        "  Other (small ops, 128 steps):  ~40%\n\n"
        "Why the FFN share drops from 68% → 33%: prefill processes 1024 "
        "tokens through one fused MoE call, so the GEMM is large and "
        "dominates. Decode does 128 separate forward passes, each on a "
        "single token. The FFN kernel call shrinks relative to the per-"
        "step overhead (router, dispatch, attention, cache write, norms) "
        "that runs on every step regardless of token count.\n\n"
        "Note the kernel naming: prefill uses cutlass_80_tensorop_bf16 "
        "(standard tensor-core GEMM optimized for large batch shapes); "
        "decode uses cutlass_80_wmma_tensorop_bf16 (WMMA variant, "
        "optimized for the small batch=1 shapes decode produces). Same "
        "CUTLASS family, different shape-specialization.\n\n"
        "Implication for FP8: still helps decode (the 14% projection "
        "GEMMs are bandwidth-bound on weight loads), and that's the "
        "bigger fraction of decode latency than people expect because "
        "the FFN kernel itself shrinks. This is consistent with the "
        "TPOT measurements: FP8 cuts ~17%, which is roughly proportional "
        "to the bandwidth-bound share of decode."
    )


def slide_lever_board(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "What I tried — and what each one bought us",
              eyebrow="LEVERS")

    fp = numbers.fp8_vs_bf16()
    pc = numbers.prefix_caching_summary()
    bd = numbers.batching_decode_summary()
    tp = numbers.tp2_vs_baseline()

    fp_cell = "TBD"
    if fp:
        L_pick = 1024 if 1024 in fp else next(iter(sorted(fp)))
        d = fp[L_pick]
        fp_cell = (f"end-to-end {abs(d['e2e_pct']):.0f}% faster, "
                   f"TPOT {abs(d['tpot_pct']):.0f}% faster")
    pc_cell = "TBD"
    if pc:
        P_max = max(pc)
        pc_cell = (f"TTFT {pc[P_max]['savings_pct']:.0f}% faster "
                   f"on shared prompts (P=2048)")
    tp_cell = "TBD"
    if tp:
        L_pick = 1024 if 1024 in tp else next(iter(sorted(tp)))
        d = tp[L_pick]
        tp_cell = (f"TPOT {abs(d['tpot_pct']):.0f}% faster; TTFT roughly flat "
                   f"at chat-app prompt sizes (L ≥ 512)")
    # Batching: honest framing — TPOT degrades, throughput rises
    bd_cell = "per-request TPOT gets worse"
    if bd:
        bd_cell = (f"per-request TPOT 4× worse at B=32; not a latency lever — "
                   f"throughput trade only")

    rows = [
        ["Change", "Status", "What it bought us (per-request)"],
        ["FP8 quantization",              "measured", fp_cell],
        ["Prefix caching",                "measured", pc_cell],
        ["Tensor parallelism (2 GPUs)",   "measured", tp_cell],
        ["Continuous batching",           "measured", bd_cell],
    ]
    add_table(s, rows, left=0.6, top=2.0, width=12.1, height=2.8, pt=15,
              col_widths=[4.2, 2.0, 5.9])
    add_caption(s,
        "Each change measured the same way against the same baseline.",
        top=5.2)
    add_footer(s, 7)
    set_notes(s,
        "Five real levers, one negative result. Each row is a script in the "
        "repo, a JSON file in results/, and (for the big ones) its own slide "
        "coming up.\n\n"
        "For the VP: this is what 'data-driven' means in practice. We did "
        "not pick the optimization that sounded cool; we picked the ones "
        "the profile told us would actually move the number we care about."
    )


def _verdict_table_only(prs: Presentation, *, title: str, eyebrow: str,
                         combo_label: str, combo_data: dict[int, dict],
                         page_num: int, notes: str) -> None:
    """Verdict table slide: full-width table, methodology footnote.

    Same shape as slide 5: rows = input length, cols = TTFT/TPOT/e2e for
    each configuration. Cells flagged red where they miss the SLA.
    """
    s = blank(prs)
    add_title(s, title, eyebrow=eyebrow)

    SLA_TTFT, SLA_TPOT, SLA_E2E = 15.0, 2.5, 500.0
    bf = numbers.derive_tpot(numbers.latency_grid("baseline"))

    rows = [
        ["Input length",
         "baseline", "baseline", "baseline",
         combo_label, combo_label, combo_label],
        ["",
         "TTFT", "TPOT", "e2e",
         "TTFT", "TPOT", "e2e"],
    ]
    cell_colors: dict[tuple[int, int], RGBColor] = {}
    for r_idx, L in enumerate(sorted(bf), start=2):
        b = bf[L]
        c = combo_data.get(L)
        rows.append([
            str(L),
            fmt_ms(b["ttft_mean_ms"]),
            f"{b['tpot_mean_ms']:.2f}",
            fmt_ms(b["end_to_end_mean_ms"], 0),
            fmt_ms(c["ttft_mean_ms"]) if c else "—",
            f"{c['tpot_mean_ms']:.2f}" if c else "—",
            fmt_ms(c["end_to_end_mean_ms"], 0) if c else "—",
        ])
        if b["ttft_mean_ms"] >= SLA_TTFT: cell_colors[(r_idx, 1)] = WARN
        if b["tpot_mean_ms"] >= SLA_TPOT: cell_colors[(r_idx, 2)] = WARN
        if b["end_to_end_mean_ms"] >= SLA_E2E: cell_colors[(r_idx, 3)] = WARN
        if c:
            if c["ttft_mean_ms"] >= SLA_TTFT: cell_colors[(r_idx, 4)] = WARN
            if c["tpot_mean_ms"] >= SLA_TPOT: cell_colors[(r_idx, 5)] = WARN
            if c["end_to_end_mean_ms"] >= SLA_E2E: cell_colors[(r_idx, 6)] = WARN

    add_table(s, rows, left=0.4, top=2.0, width=12.5, height=3.6, pt=15,
              col_widths=[1.6, 1.6, 1.4, 1.7, 1.6, 1.4, 3.2],
              cell_text_colors=cell_colors)
    add_text(s,
             "SLA: TTFT < 15 ms, TPOT < 2.5 ms (mean), e2e < 500 ms (200-token "
             "response). Red cells miss the SLA.",
             left=0.6, top=5.8, width=12.0, height=0.4,
             pt=12, italic=True, color=SUBTLE)
    add_text(s,
             "Baseline numbers from in-process vLLM bench. Caching combo from "
             "HTTP harness, with ~20 ms harness overhead reconciled per input "
             "length against the bf16 cache-off harness baseline.",
             left=0.6, top=6.25, width=12.0, height=0.6,
             pt=10, italic=True, color=SUBTLE)
    add_footer(s, page_num)
    set_notes(s, notes)


def _verdict_chart_only(prs: Presentation, *, title: str, eyebrow: str,
                         chart_path: Path, page_num: int, notes: str) -> None:
    """Verdict chart slide: full-width end-to-end bar chart."""
    s = blank(prs)
    add_title(s, title, eyebrow=eyebrow)
    add_picture(s, chart_path,
                left=1.0, top=1.95, width=11.3, height=4.85)
    add_footer(s, page_num)
    set_notes(s, notes)


def slide_sla_verdict_table_recommended(prs: Presentation) -> None:
    combo = numbers.reconciled_caching_combo("prefix_caching_fp8") or {}
    _verdict_table_only(
        prs,
        title="SLA verdict — FP8 + prefix caching",
        eyebrow="VERDICT",
        combo_label="FP8 + prefix caching",
        combo_data=combo,
        page_num=8,
        notes=(
            "Baseline on the left, FP8 + prefix caching on the right. "
            "Red cells miss the SLA.\n\n"
            "Baseline starts going red at L=512 (e2e fails) and is fully "
            "red at L=1024 and L=1500 (all three metrics fail).\n\n"
            "FP8 + prefix caching is green across every chat-app input "
            "length. TTFT compresses to ~7-8 ms (from 10-25 ms). TPOT "
            "trims to 2.0-2.4 ms. End-to-end ranges from 396 ms at L=128 "
            "to 486 ms at L=1500 — comfortably under the 500 ms target "
            "across the whole grid.\n\n"
            "Methodology: caching combo measured via HTTP harness, which "
            "adds ~20 ms overhead per request. Reconciled per-L against "
            "the bf16 cache-off harness baseline. Per-L overhead is "
            "consistent at 19-21 ms — subtraction is stable."
        ),
    )


def slide_sla_verdict_chart_ttft(prs: Presentation) -> None:
    _verdict_chart_only(
        prs,
        title="SLA verdict — TTFT across input length",
        eyebrow="VERDICT",
        chart_path=FIG_DIR / "appendix_verdict_metric_ttft.png",
        page_num=10,
        notes=(
            "TTFT verdict, all three configurations side by side at every "
            "chat-app input length.\n\n"
            "Grey bars: baseline. Cross the 15 ms threshold at L=1024 "
            "(19.6 ms) and L=1500 (25.1 ms) — those bars turn red.\n\n"
            "Green bars: FP8 + prefix caching. ~7-8 ms across the entire "
            "grid — well under the 15 ms target. Caching collapses TTFT "
            "to roughly the suffix cost regardless of total prompt length, "
            "so the line is essentially flat.\n\n"
            "Blue bars: FP8 + prefix caching + TP=2. Tracks the green bars "
            "almost identically — TP=2 trims another ~0.5-1 ms but the "
            "win is below production noise.\n\n"
            "Headline: FP8 + prefix caching alone clears the TTFT bar at "
            "every chat-app prompt size with margin."
        ),
    )


def slide_sla_verdict_chart_tpot(prs: Presentation) -> None:
    _verdict_chart_only(
        prs,
        title="SLA verdict — TPOT across input length",
        eyebrow="VERDICT",
        chart_path=FIG_DIR / "appendix_verdict_metric_tpot.png",
        page_num=11,
        notes=(
            "TPOT verdict. Threshold is 2.5 ms per output token.\n\n"
            "Grey bars (baseline): cross the line at L=512 (2.47 ms) and "
            "stay over through L=1024 and L=1500. Red bars at L≥512.\n\n"
            "Green bars (FP8 + prefix caching): under the threshold at "
            "every input length. FP8 is the lever that moves TPOT — it "
            "halves the per-token weight-load bandwidth demand. Caching "
            "doesn't affect TPOT (only TTFT), but the green bars below "
            "show TPOT is comfortably inside SLA.\n\n"
            "Blue bars (everything on): essentially identical to green. "
            "TP=2 trims another ~0.02 ms per token — invisible at this "
            "scale.\n\n"
            "Headline: TPOT is the easier of the two metrics to clear. "
            "FP8 alone would do it; caching and TP=2 are bonus."
        ),
    )


def slide_sla_verdict_chart_e2e(prs: Presentation) -> None:
    _verdict_chart_only(
        prs,
        title="SLA verdict — end-to-end across input length",
        eyebrow="VERDICT",
        chart_path=FIG_DIR / "appendix_verdict_metric_e2e.png",
        page_num=12,
        notes=(
            "End-to-end verdict — the headline metric. Threshold is "
            "500 ms for a 200-token response.\n\n"
            "Grey bars (baseline): miss the SLA from L=512 onward "
            "(505, 565, 603 ms). Red bars at L≥512.\n\n"
            "Green bars (FP8 + prefix caching): every bar under 500 ms. "
            "Range 396 ms (L=128) to 486 ms (L=1500) — comfortably "
            "inside the SLA across the whole chat-app range. Margin "
            "shrinks as input grows but never crosses.\n\n"
            "Blue bars (everything on): 392 to 483 ms. ~3-5 ms better "
            "than green at every L. Real but marginal.\n\n"
            "Headline: FP8 + prefix caching is the recommended deployment. "
            "Everything-on is a backup option if TPOT or memory becomes "
            "a problem at higher concurrency — not for this single-user "
            "SLA story."
        ),
    )


def slide_sla_verdict_table_all_on(prs: Presentation) -> None:
    combo = numbers.reconciled_caching_combo("prefix_caching_fp8_tp2") or {}
    _verdict_table_only(
        prs,
        title="SLA verdict — everything on (+ TP=2)",
        eyebrow="VERDICT",
        combo_label="FP8 + caching + TP=2",
        combo_data=combo,
        page_num=9,
        notes=(
            "Same baseline, but the combo column is everything-on: FP8 + "
            "prefix caching + TP=2 across both Blackwell GPUs.\n\n"
            "All cells green, same as the previous verdict. TTFT shaves "
            "another ~1 ms off the FP8+caching numbers; TPOT trims "
            "another ~0.02 ms; end-to-end is ~3-5 ms better. TP=2 buys "
            "real but marginal additional latency on top of "
            "FP8+caching.\n\n"
            "Recommendation stays: FP8 + prefix caching in prod. "
            "Evaluate TP=2 if TPOT becomes the bottleneck under higher "
            "concurrency — but for the chat-app per-request SLA, the "
            "two-lever combo is enough."
        ),
    )




def slide_tpot_lines(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "TPOT — baseline vs two levers", eyebrow="RESULTS")
    add_picture(s, FIG_DIR / "appendix_tpot_vs_length.png",
                left=1.4, top=1.95, width=10.5, height=4.85)
    add_footer(s, 9)
    set_notes(s,
        "TPOT counterpart to slide 8. Same three lines, different metric — "
        "three of the four levers (bf16 baseline, FP8, TP=2). Continuous "
        "batching is on slide 12 on its own axis (batch size, not prompt "
        "length), because it's not a per-request latency lever — it's a "
        "throughput lever that makes per-request TPOT worse.\n\n"
        "Blue line — bf16 baseline. Roughly flat across input length: 2.34 "
        "ms/token at L=128 up to 2.90 ms/token at L=1500. Decode TPOT barely "
        "depends on prompt length because each decode step does the same "
        "weight-load work regardless of how long the cached context is.\n\n"
        "Green line — FP8. Consistent ~15–17% TPOT improvement across all "
        "prompt sizes. FP8 attacks the bandwidth-bound weight loads that "
        "dominate decode, so the savings are stable.\n\n"
        "Red line — TP=2. ~12–14% TPOT gain across the grid. Splitting the "
        "model across both GPUs halves the per-GPU weight-load bandwidth "
        "demand per token, so the bandwidth-bound decode loop runs faster.\n\n"
        "The takeaway: for TPOT — unlike TTFT — both FP8 and TP=2 deliver "
        "real, consistent wins. They're additive in principle, though I "
        "didn't measure the combination."
    )


def slide_big_wins(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "TTFT — baseline vs two levers", eyebrow="RESULTS")
    add_picture(s, FIG_DIR / "appendix_ttft_vs_length.png",
                left=1.4, top=1.95, width=10.5, height=4.85)
    add_footer(s, 8)
    set_notes(s,
        "One chart, three of the four levers (bf16 baseline, FP8, TP=2). "
        "Continuous batching is on slide 12 on its own axis (batch size, "
        "not prompt length), because it's not a per-request latency lever — "
        "it's a throughput lever that makes per-request latency worse.\n\n"
        "Walk the audience across the three lines on this chart:\n\n"
        "Blue line — bf16 baseline. The shape of the problem.\n\n"
        "Green line — FP8 quantization. Consistently below baseline by "
        "15–35%. It's a flag flip on vLLM; we'll cover the accuracy "
        "tradeoff on the next slide.\n\n"
        "Red line — TP=2 across both Blackwell GPUs. Slightly better at "
        "L=128, sits on top of the baseline at higher L. The reason is "
        "subtle: at batch=1 the per-step work is small, so the cross-GPU "
        "all-reduce dominates whatever parallelism we gain. TP=2 is a "
        "throughput lever more than a TTFT lever.\n\n"
        "The takeaway: for this chat workload, FP8 is the lever that moves "
        "TTFT. TP=2 helps when the GPU has more work per step — which is "
        "exactly the regime we'll hit when we start batching."
    )


def slide_tradeoffs(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "Prefix caching — payoff scales with shared-prefix length",
              eyebrow="TRADE-OFFS")
    add_picture(s, FIG_DIR / "appendix_prefix_caching.png",
                left=1.4, top=1.95, width=10.5, height=4.85)
    add_footer(s, 10)
    set_notes(s,
        "Our chat app uses a long system prompt with tool descriptions and "
        "few-shot examples — that's the right side of this chart, where the "
        "savings are largest.\n\n"
        "Left bars: short shared prompt, 13% faster first word. "
        "Middle: medium prompt, 14%. Right: long shared prompt, 40% faster.\n\n"
        "The pattern: reuse helps most when the prompt is most expensive to "
        "recompute — which is exactly when we want to skip it.\n\n"
        "Caveat for a different workload: if every request has a totally "
        "unique prompt (e.g. RAG with a different retrieved document each "
        "time), this lever does nothing.\n\n"
        "Also worth flagging if asked: lower-precision math (FP8) has its "
        "own trade-off. Three out of four standard quality benchmarks are "
        "unchanged within noise. The fourth, grade-school math, drops about "
        "2.8 percentage points. We'd test on the customer's own evaluations "
        "before flipping that switch in production."
    )


def slide_business_impact(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "Continuous batching is a throughput lever, not a latency one",
              eyebrow="APPENDIX")
    add_picture(s, FIG_DIR / "appendix_batching_pareto.png",
                left=1.4, top=1.95, width=10.5, height=4.85)
    add_footer(s, 20)
    set_notes(s,
        "Honest framing of the fourth lever. The bars rise — total throughput "
        "goes up nearly 14× from B=1 to B=128. The line also rises — per-"
        "request TPOT goes up 4× by B=32 and 4.6× by B=128.\n\n"
        "So under a per-request latency frame, batching is not a win — it's "
        "a deliberate trade. You spend per-request latency to buy total "
        "throughput per pod.\n\n"
        "When that trade is worth it: the workload is throughput-sensitive "
        "rather than latency-sensitive. Batch document processing, async "
        "summarization, anything where total tokens served per dollar matters "
        "more than how fast any single response streams.\n\n"
        "When it isn't: interactive chat. The whole point of chat is "
        "responsiveness. Don't batch past the point where TPOT crosses the "
        "user's tolerance — for a 30 ms target, B=1 is fine, B=32 is fine "
        "(9.6 ms), but you're spending real per-request latency for the "
        "throughput.\n\n"
        "If the panel asks: yes, batching helps cost per token. The chart "
        "is here so the choice is informed, not so it's automatic."
    )


def slide_sla_status(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "SLA status — where the lines hold and where they break",
              eyebrow="SLA")
    add_picture(s, FIG_DIR / "appendix_sla_status.png",
                left=1.4, top=1.95, width=10.5, height=4.85)
    add_footer(s, 12)
    set_notes(s,
        "Direct answer to 'did we meet the SLA': yes, up to batch=32 "
        "in-flight requests per pod. Both bf16 and FP8 hold under the "
        "500 ms TTFT p90 line at B=1, 4, 16, 32. Both fail at B=64 and "
        "B=128.\n\n"
        "FP8 doesn't meaningfully extend the SLA-passing concurrency "
        "ceiling. It improves TPOT and per-batch throughput, but the "
        "binding constraint at high concurrency is queueing, not the "
        "kernel — so faster math doesn't help when the engine is "
        "already saturated. Note FP8's TTFT p90 actually trails bf16 at "
        "B=64 and B=128; both are well past the SLA at that point so the "
        "ordering is noise on top of an overloaded engine.\n\n"
        "TPOT (not shown on this chart, but in the appendix) passes the "
        "30 ms target at every batch size for both — TPOT is a non-binding "
        "constraint here. The whole SLA story comes down to TTFT under "
        "concurrent load.\n\n"
        "What this means: the levers don't *rescue* the SLA — the baseline "
        "already passed it at expected load. The levers extend per-request "
        "headroom and cut cost per request. To extend the concurrency "
        "ceiling specifically, the next step is horizontal scaling (more "
        "pods), not faster kernels."
    )


def slide_roi(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "What this is worth", eyebrow="BUSINESS VALUE")

    # Headline callout — the answer to the title slide question
    callout = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.6), Inches(1.85),
                                 Inches(12.1), Inches(0.7))
    callout.fill.solid()
    callout.fill.fore_color.rgb = ACCENT
    callout.line.fill.background()
    tf = callout.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.text = "Production launch is unblocked"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        r.font.name = FONT
        r.font.size = Pt(24)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    # Vertically center the text in the box
    callout.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Left column — per-request speed (measured: 565 → 455 ms at L=1024)
    add_text(s, "PER-REQUEST SPEED", left=0.6, top=2.85, width=6.0, height=0.3,
             pt=11, bold=True, color=ACCENT)
    add_text(s, "~19% faster on every chat-app request",
             left=0.6, top=3.20, width=6.2, height=0.6,
             pt=22, bold=True, color=INK)
    add_text(s,
             "End-to-end at L=1024 drops from 565 ms to 455 ms with FP8 + "
             "prefix caching — measured. 110 ms shaved off every request "
             "across the chat-app prompt range.",
             left=0.6, top=3.95, width=6.0, height=1.4,
             pt=13, color=SUBTLE)
    add_text(s, "Break-even: Day 1",
             left=0.6, top=5.55, width=6.0, height=0.5,
             pt=20, bold=True, color=INK)
    add_text(s,
             "Both levers are config changes — no new hardware, no "
             "downtime, no model retraining.",
             left=0.6, top=6.10, width=6.0, height=0.9,
             pt=13, color=SUBTLE)

    # Right column — cost per request (g7e.12xlarge, $8.29/hr)
    add_text(s, "COST PER REQUEST", left=7.0, top=2.85, width=6.0, height=0.3,
             pt=11, bold=True, color=ACCENT)
    add_text(s, "$1.30 → $1.05 per 1K requests",
             left=7.0, top=3.20, width=6.0, height=0.6,
             pt=22, bold=True, color=INK)
    add_text(s,
             "On AWS g7e.12xlarge ($8.29 / hr, 2× RTX PRO 6000 Blackwell — "
             "the same GPU we measured), FP8 + prefix caching cuts cost per "
             "request 19%. At 1M requests / day, that saves ~$93K per pod "
             "per year.",
             left=7.0, top=3.95, width=6.0, height=1.6,
             pt=13, color=SUBTLE)
    add_text(s, "Same hardware, same code path",
             left=7.0, top=5.55, width=6.0, height=0.5,
             pt=20, bold=True, color=INK)
    add_text(s,
             "No re-architecture, no new dependencies, no ops change. "
             "Re-baseline whenever the model, vLLM version, or hardware "
             "changes.",
             left=7.0, top=6.10, width=6.0, height=0.9,
             pt=13, color=SUBTLE)

    add_footer(s, 13)
    set_notes(s,
        "Two stories for the VP. Per-request speed on the left, cost per "
        "request on the right. Both backed by directly measured FP8 + "
        "prefix caching numbers (verdict slides 12-15).\n\n"
        "Speed: end-to-end at L=1024 drops from 565 ms to 455 ms with FP8 "
        "+ prefix caching — that's 110 ms shaved off every request, 19% "
        "faster. Measured directly, not estimated.\n\n"
        "Cost math: AWS g7e.12xlarge is $8.29 per hour on-demand. That's "
        "the instance with 2× RTX PRO 6000 Blackwell — exactly the GPU "
        "this work was measured on, so we're not proxying with a "
        "different SKU. $8.29 / hr divided by 3600 seconds is $0.00230 "
        "per second of GPU time. A baseline 565 ms request costs about "
        "$1.30 per thousand. A 455 ms FP8+caching request costs about "
        "$1.05 per thousand. That's 25 cents saved per thousand requests, "
        "or about $254 per million. At 1 million requests per day — a "
        "reasonable production chat-app rate — that's $254/day, $93,000 "
        "per pod per year.\n\n"
        "Both are config changes. No model retraining, no new "
        "dependencies, no ops change. Break-even on day one. The only "
        "ongoing discipline is to re-baseline whenever the model, vLLM "
        "version, or hardware moves — the harness is set up to make that "
        "cheap.\n\n"
        "If pressed on the dollar number: the rate is AWS public list "
        "price; an enterprise deployment with reserved capacity would see "
        "smaller dollar savings but the same percentage. The $93K is "
        "linear in traffic, so 100K requests/day = $9.3K/year, 10M "
        "requests/day = $930K/year."
    )


def slide_architecture_tradeoffs(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "Trade-offs — nothing is free", eyebrow="TRADE-OFFS")
    rows = [
        ["Technique", "Upside", "Risk", "Mitigation"],
        ["FP8 quantization",
         "17% faster end-to-end",
         "Slight precision loss on some ops",
         "Run eval suite before prod; bf16 fallback available"],
        ["Prefix caching",
         "40% TTFT cut on shared prompts",
         "Cache miss = no benefit; memory overhead",
         "Size KV cache to workload; monitor hit rate"],
        ["Tensor parallel (2 GPUs)",
         "14% TPOT gain at single-user load",
         "TTFT win is small at chat-app prompt sizes (L ≥ 512); concurrency behavior not measured",
         "Use for TPOT bottlenecks; re-measure before scaling under load"],
        ["Continuous batching",
         "Up to 14× pod throughput (cost-per-token)",
         "Per-request TPOT ~4× worse at B=32 — not a latency lever",
         "Use only when throughput beats latency — batch jobs, not chat"],
    ]
    add_table(s, rows, left=0.4, top=2.0, width=12.5, height=4.4, pt=12,
              col_widths=[2.4, 2.6, 3.7, 3.8])
    add_footer(s, 14)
    set_notes(s,
        "Every lever has a price; this is the slide where I name them so "
        "neither panelist has to dig them out of me. Walk the table top-to-"
        "bottom.\n\n"
        "FP8: the upside is real, the risk is GSM8K-style precision-sensitive "
        "tasks taking a measurable hit. The mitigation is to run the eval "
        "suite on the customer's actual workload before flipping the flag in "
        "production, and to keep bf16 as a fallback if anything regresses.\n\n"
        "Prefix caching: only helps when prompts share prefixes. Sizing the "
        "KV cache and monitoring hit rate keeps it honest in production.\n\n"
        "Tensor parallel: I keep flagging this — at single-user batch=1 it "
        "doesn't move TTFT. The reason it's still worth it is TPOT and the "
        "doubled memory headroom for bigger batches. If someone uses it as a "
        "TTFT lever they'll be disappointed.\n\n"
        "Concurrency: throughput goes up, per-user latency goes up. The "
        "mitigation is to set max_num_seqs from the p90 TTFT we promised, "
        "not from the throughput we want."
    )


def slide_scaling_10x(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "If traffic grows 10× tomorrow", eyebrow="SCALE")
    rows = [
        ["#", "Action", "Why now", "Breaks when"],
        ["1",
         "Enable FP8 + prefix caching",
         "Free, today — config only",
         "Accuracy regresses on customer evals, or prompts don't share prefixes"],
        ["2",
         "Tune concurrency to the SLA",
         "Same pod, no new hardware",
         "Sustained load exceeds ~200–600 active users per pod (batch=32 in flight)"],
        ["3",
         "Use both GPUs together (TP=2)",
         "Same node, +14% TPOT, 2× memory",
         "One pod still saturated — need more pods, not bigger ones"],
        ["4",
         "Horizontal pod autoscaling (HPA)",
         "Linear capacity, k8s-native",
         "GPU node pool runs out of room"],
        ["5",
         "Add new GPU machines",
         "Capacity beyond the current cluster",
         "Last resort — triggers procurement cycle"],
    ]
    add_table(s, rows, left=0.4, top=2.0, width=12.5, height=4.6, pt=12,
              col_widths=[0.5, 3.4, 3.6, 5.0])
    add_footer(s, 15)
    set_notes(s,
        "Standard scaling playbook for the EKS shape we're in. Each row is "
        "cheaper and lower-risk than the next; only move down a row when the "
        "row above stops being enough.\n\n"
        "Row 1 — software config. FP8 plus prompt reuse is free. The only "
        "way it stops working is if FP8 hurts accuracy on the customer's own "
        "evals, or if traffic doesn't actually share prefixes.\n\n"
        "Row 2 — concurrency tuning. Set max_num_seqs from the p90 TTFT "
        "target, not from peak throughput. The measured ceiling is "
        "batch=32 in flight per pod — past that, TTFT p90 exceeds the "
        "500 ms target. Translating that to humans: at 5–15% in-flight "
        "rate (typical chat with think-time), one pod can absorb roughly "
        "200–600 logged-in active users before this lever runs out. "
        "Caveat: that range depends on think-time, prompt length, and "
        "session shape — not measured against the actual chat-app traffic "
        "yet. When this row stops being enough, add pods (row 4) before "
        "doing anything bigger.\n\n"
        "Row 3 — TP=2 inside one pod. Uses both Blackwell GPUs together. "
        "It buys ~14% TPOT and doubles the memory available for bigger "
        "batches. It does NOT buy TTFT at single-user load. When this row "
        "stops being enough we add horizontal capacity, not vertical.\n\n"
        "Row 4 — horizontal pod autoscaling. Add more replicas behind the "
        "load balancer. Linear capacity scaling, k8s-native, no engineer "
        "in the loop. Stops when the GPU node pool is full.\n\n"
        "Row 5 — new nodes. Most expensive lever; triggers the procurement "
        "cycle. We push it as far right as possible."
    )


def slide_recommendation(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "Recommendation", eyebrow="VERDICT")
    add_bullets(s, [
        "Enable FP8 quantization and prefix caching this week — "
        "validate accuracy on customer evals first",
        "Evaluate tensor parallelism (TP=2) for TPOT — 14% gain at "
        "single-user load; rerun under concurrency before adopting in prod",
        "Use continuous batching only for throughput-mode workloads — "
        "not chat — per-request TPOT degrades 4× by B=32",
        "Re-baseline whenever the model, vLLM version, or hardware changes",
    ])
    add_footer(s, 15)
    set_notes(s,
        "Five lines, in priority order. A staff engineer reading this "
        "should agree with each one.\n\n"
        "Last bullet matters: every claim here is tied to *this* model, "
        "*this* engine version, *this* hardware. Re-baseline whenever "
        "any of those move. The repo is set up to make that cheap."
    )


def slide_scaling_remeasure(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "When scale increases — re-measure",
              eyebrow="WHAT'S NEXT")

    rows = [
        ["Scale increase", "What to expect", "Action"],
        ["Input length grows toward 4K context",
         "TTFT scales ~linearly with input length; SLA margin shrinks",
         "Re-run benchmark_latency.py at the new L; verify SLA holds"],
        ["Larger model — more experts, similar active params",
         "Total weights and KV cache grow; per-token TPOT roughly unchanged",
         "Re-measure memory headroom; TP=2 may become required"],
    ]
    add_table(s, rows, left=0.4, top=2.0, width=12.5, height=2.4, pt=14,
              col_widths=[4.0, 4.4, 4.1])
    add_text(s,
             "Every row in this deck is tied to this model + this engine + "
             "this hardware. The harness (`scripts/`) makes re-measurement a "
             "one-command rerun — re-baseline before committing to numbers "
             "that aren't on a current deck.",
             left=0.4, top=6.20, width=12.5, height=0.7,
             pt=12, italic=True, color=SUBTLE)
    add_footer(s, 16)
    set_notes(s,
        "What to do when the workload changes. Frames input-length growth, "
        "larger-model swaps, and stack updates as triggers to re-measure — "
        "not as numerical projections.\n\n"
        "Why no projection numbers: anything I quote past the measured grid "
        "is extrapolation. The honest answer is 'rerun the benchmark in 30 "
        "minutes when the workload changes' — and the harness is set up to "
        "make that cheap.\n\n"
        "If the panel pushes for a number: I can sketch first-principles "
        "scaling — TTFT scales linearly with input length at fixed "
        "model size, TPOT scales with active-parameter bytes loaded per "
        "token. But the slide deliberately avoids putting projected numbers "
        "in front of the panel because they wouldn't be defensible without "
        "actual data.\n\n"
        "Specifically:\n"
        "- Input length toward 4K: linear extrapolation suggests baseline "
        "fails harder, FP8+caching gets borderline by ~3000 tokens.\n"
        "- 10× active params: not recoverable on single Blackwell — needs "
        "multi-GPU TP=4/8 or different hardware class.\n"
        "- More experts at same active: KV cache and weights scale; TP=2 "
        "becomes mandatory rather than optional."
    )


def _add_code_block(slide, lines: list[str], *, left: float, top: float,
                     width: float, height: float, pt: int = 11) -> None:
    """Render lines as monospace text in a code-block style.

    Comments (anything after '#') are colored SUBTLE so the file/dir names
    pop. Indentation is preserved; matplotlib-style trees render correctly.
    """
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(0)

        # Split into "code" and "# comment" parts
        if "#" in line:
            code, _, comment = line.partition("#")
            r1 = p.add_run()
            r1.text = code
            r1.font.name = "Courier New"
            r1.font.size = Pt(pt)
            r1.font.color.rgb = INK
            r2 = p.add_run()
            r2.text = "#" + comment
            r2.font.name = "Courier New"
            r2.font.size = Pt(pt)
            r2.font.color.rgb = SUBTLE
        else:
            r = p.add_run()
            r.text = line
            r.font.name = "Courier New"
            r.font.size = Pt(pt)
            r.font.color.rgb = INK


def slide_code_overview(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "What I wrote", eyebrow="THE CODE")

    # scripts/
    add_text(s, "scripts/", left=0.5, top=1.85, width=12.0, height=0.3,
             pt=14, bold=True, color=ACCENT)
    _add_code_block(s, [
        "  smoke_test.py              # sanity check before benchmarking — loads OLMoE, generates one token",
        "  benchmark_latency.py       # direct-from-engine TTFT/TPOT (no HTTP overhead) — bf16, --fp8, --tp 2",
        "  profile_prefill.py         # torch.profiler trace at L=1024 — kernel breakdown",
        "  profile_decode.py          # torch.profiler trace at I=1, O=128 — decode kernels",
        "  bench_prefix_caching.py    # caching via HTTP harness (bf16 / fp8 / tp2 combos)",
        "  bench_batching_decode.py   # tokens/sec and per-user TPOT across concurrency 1 → 128",
        "  eval_suite.py              # lm-evaluation-harness — FP8 accuracy on MMLU / GSM8K / ...",
        "  harness/                   # engine-agnostic vLLM-server launcher + HTTP client",
    ], left=0.5, top=2.25, width=12.5, height=3.0, pt=13)

    # results/
    add_text(s, "results/", left=0.5, top=5.40, width=12.0, height=0.3,
             pt=14, bold=True, color=ACCENT)
    _add_code_block(s, [
        "  baseline_*.json            # bf16 / FP8 / TP=2 latency cells (one JSON per L × O)",
        "  harness/                   # HTTP-harness outputs (prefix caching, batching)",
        "  profile/                   # torch.profiler traces + per-kernel summaries",
        "  eval/                      # lm-evaluation-harness raw outputs",
    ], left=0.5, top=5.80, width=12.5, height=1.1, pt=13)

    add_footer(s, 18)
    set_notes(s,
        "One-slide summary of the codebase, right before Q&A and the repo "
        "link. Reads like the README — actual file tree, file names tell "
        "the story.\n\n"
        "TWO MEASUREMENT PATHS, TWO DIFFERENT JOBS:\n\n"
        "1) In-process — `vllm bench latency` driven by "
        "scripts/benchmark_latency.py and the profile_*.py scripts. The "
        "bench tool spins up vLLM in the same Python process, fires "
        "synthetic prompts directly through the engine API, and times "
        "from inside the process. No HTTP, no network, no scheduler "
        "queue. This is what you want when you're measuring the engine "
        "itself: TTFT/TPOT for the kernels in isolation, with the "
        "tightest variance and no harness overhead. Every number on "
        "slides 5, 8, 9 (lever-only rows) comes from this path.\n\n"
        "2) HTTP harness — scripts/harness/ launches a real `vllm serve` "
        "subprocess, waits for the OpenAI-compatible /v1/models endpoint, "
        "issues streaming HTTP requests, and times from the client side. "
        "Every number on slides 8, 9 (caching combo rows), 10–12, and the "
        "appendix batching slide comes from this path. Adds about 20 ms "
        "of overhead per request vs the in-process bench — measured and "
        "reconciled per input length (verdict slide methodology footnote).\n\n"
        "WHY WE NEED A HARNESS:\n\n"
        "Three things the in-process bench can't measure, but the harness "
        "can:\n\n"
        "Prefix caching. The bench fires unrelated random-token prompts "
        "with no shared prefix between requests. There's nothing for the "
        "cache to hit on. The harness builds a workload of 30 requests "
        "that all share a fixed prefix and vary only a 64-token suffix — "
        "exactly what a chat app with a shared system prompt looks like. "
        "Without this shape, prefix caching can't be measured at all.\n\n"
        "Concurrent load. The bench runs requests sequentially. The "
        "harness drives a configurable number of in-flight requests "
        "against the live server, the same way real production traffic "
        "hits the engine. That's how we know batch=32 saturates the "
        "engine and batch=64 starts queueing.\n\n"
        "Real serving path. `vllm serve` is what production runs. The "
        "in-process bench skips the API layer, the scheduler, the "
        "request queue, and the streaming response code. The harness "
        "exercises all of that, so the numbers reflect what users will "
        "actually feel.\n\n"
        "REALISTIC PROMPTS — SHAREGPT: bench_batching.py loads "
        "real chat prompts from the ShareGPT V3 dataset "
        "(huggingface.co/datasets/anon8231489123/ShareGPT_Vicuna_unfiltered) "
        "via scripts/harness/datasets.py. Prompts are filtered to a "
        "128–512 token band and pooled so concurrent requests don't "
        "accidentally share a prefix (which would let prefix caching "
        "merge them and inflate throughput). Synthetic random tokens "
        "would still produce numbers, but ShareGPT's distribution of "
        "real questions, code snippets, and multi-turn snippets is "
        "closer to what production traffic looks like — so when the "
        "throughput chart says 'X tok/s at concurrency Y,' that holds "
        "for chat-shaped traffic, not just synthetic stress. The pure-"
        "decode counterpart (bench_batching_decode.py) keeps prompts "
        "synthetic on purpose because it isolates MoE-batching behavior "
        "from prefill-mixing.\n\n"
        "ENGINE-AGNOSTIC: the harness wraps server lifecycle behind an "
        "OpenAI-compatible HTTP client. Swap the launcher and the same "
        "benchmarks run against TRT-LLM, SGLang, anything with an "
        "OpenAI endpoint. Deliberate — when the next model arrives, the "
        "harness doesn't have to be rewritten.\n\n"
        "REPRODUCIBILITY: every measurement → timestamped JSON in "
        "results/. talk/numbers.py reads JSON; deck refreshes on rerun. "
        "scripts/plots/ regenerates every chart from results/. Seed "
        "pinned at 42 — variance across runs is under 1 ms, which is "
        "how we trust small deltas as signal, not noise."
    )


def slide_repo(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "Code, results, and this deck",
              eyebrow="REPO")

    # QR centered
    add_picture(s, FIG_DIR / "repo_qr.png",
                left=5.0, top=2.0, width=3.5, height=3.5)
    add_text(s, "github.com/harishvs/moe-inference-optimize",
             left=2.0, top=5.6, width=9.3, height=0.5,
             pt=20, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(s, "MIT license — clone, run, contribute",
             left=2.0, top=6.15, width=9.3, height=0.4,
             pt=13, italic=True, color=SUBTLE, align=PP_ALIGN.CENTER)
    add_footer(s, 17)
    set_notes(s,
        "QR slide. Lingers while the panel scans. The repo includes:\n\n"
        "- Every script that produced every number in the deck\n"
        "- All raw JSON measurements under results/\n"
        "- The deck itself (build_deck.py + talk/numbers.py)\n"
        "- README walks through reproduction step by step\n"
        "- MIT-licensed; clone, fork, send PRs\n\n"
        "Stay on this slide for ~10 seconds before transitioning to Q&A "
        "so anyone scanning has time to capture the link."
    )


def slide_qa(prs: Presentation) -> None:
    s = blank(prs)
    add_text(s, "Q&A",
             left=0.8, top=2.5, width=12.0, height=2.0,
             pt=120, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_rule(s, left=5.6, top=4.6, width=2.0, height_pt=4.0)
    add_text(s, "15 minutes — fire away.",
             left=0.8, top=4.9, width=12.0, height=0.5,
             pt=22, color=SUBTLE, align=PP_ALIGN.CENTER)
    add_footer(s, 19)
    set_notes(s,
        "Anticipated questions:\n\n"
        "Q (Tech): Does FP8 still win on Blackwell when bandwidth is much "
        "higher? A: Yes — measured. The win is smaller than on L4 because "
        "Blackwell's HBM3e and tensor cores already eat much of the "
        "bandwidth bottleneck FP8 attacks.\n\n"
        "Q (Tech): What about FP4? A: Blackwell supports it; vLLM/OLMoE "
        "compatibility is the limiter. On the roadmap; not measured here.\n\n"
        "Q (Tech): Does FP8 affect MoE routing — which experts fire? A: "
        "vLLM's MoE FP8 quantization keeps the router (gating linear) in "
        "higher precision by design; only the expert FFN weights are "
        "quantized. So expert selection at any given token is identical "
        "between bf16 and FP8 — only the math each selected expert does "
        "changes precision. This is consistent with the eval results: "
        "broad routing instability would have hit MMLU/HellaSwag/ARC, not "
        "just GSM8K. GSM8K's bigger drop is consistent with low-precision "
        "arithmetic compounding errors in chain-of-thought math, not with "
        "a routing change. Caveat: I haven't directly compared expert-fire "
        "patterns trace-by-trace between bf16 and FP8 — the claim is from "
        "reading vLLM's quantization code path plus the eval shape.\n\n"
        "Q (Tech): Tensor parallel vs expert parallel for MoE? A: Out of "
        "scope for this talk. EP is the obvious next experiment — expert "
        "parallelism specifically targets the MoE structure.\n\n"
        "Q (VP): What's the dollar number on $/M tokens? A: Depends on "
        "the instance — roughly the throughput ratio at the batching knee, "
        "applied to the EKS GPU node pool's cost. Concrete dollar estimate "
        "available in a follow-up.\n\n"
        "Q (VP): What's the time horizon? A: FP8 + prefix caching this "
        "week. TP=2 next sprint. HPA tuning is ongoing.\n\n"
        "Q (VP): What if the model changes? A: The harness is engine- and "
        "model-agnostic. Re-baseline runs in ~30 minutes.\n\n"
        "Roadmap if asked 'what would another sprint buy us?': expert-parallel "
        "sharding (cross-GPU MoE), FP4 once vLLM/OLMoE wiring lands, "
        "speculative decoding with a small draft model for chat workloads, "
        "tuning to the actual customer prompt distribution, and continuous-"
        "batching policy tuning under bursty real traffic."
    )


# ---------- Appendix ----------

def appendix_fp8_accuracy(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "FP8 quality — standard accuracy benchmarks",
              eyebrow="APPENDIX")
    add_picture(s, FIG_DIR / "appendix_fp8_accuracy.png",
                left=1.4, top=1.95, width=10.5, height=4.85)
    add_footer(s, 21)


def appendix_ttft_lines(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "TTFT across input length, by lever", eyebrow="APPENDIX")
    add_picture(s, FIG_DIR / "appendix_ttft_vs_length.png",
                left=0.6, top=1.85, width=12.1)
    add_caption(s,
        "FP8 wins consistently. TP=2 mostly tracks baseline at small batch — "
        "the cross-GPU sync overhead eats the parallelism win on prefill alone.",
        top=6.5)
    add_footer(s, 15)


def appendix_batching(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "Batching Pareto — pick latency or throughput",
              eyebrow="APPENDIX")
    add_picture(s, FIG_DIR / "appendix_batching_pareto.png",
                left=0.6, top=1.85, width=12.1)
    add_caption(s,
        "B=1 for chat-latency users; near the knee (B=16–32) for cost-mode "
        "batch jobs. The latency target decides which side of this curve you operate on.",
        top=6.5)
    add_footer(s, 16)


def appendix_prefix_caching(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "Prefix caching — payoff scales with shared-prefix length",
              eyebrow="APPENDIX")
    add_picture(s, FIG_DIR / "appendix_prefix_caching.png",
                left=0.6, top=1.85, width=12.1)
    add_caption(s,
        "Workload-dependent: zero benefit if every request has a unique prompt; "
        "large benefit on shared system-prompt traffic (the chat-app case).",
        top=6.5)
    add_footer(s, 17)


def appendix_cost_math(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "Cost per request — the math",
              eyebrow="APPENDIX")

    rows = [
        ["Step",                              "Value",                    "Source"],
        ["AWS g7e.12xlarge on-demand rate",  "$8.29 / hour",             "AWS EC2 public pricing"],
        ["Per-second GPU-pod cost",           "$8.29 ÷ 3600 = $0.002303", "arithmetic"],
        ["Baseline e2e (L=1024, O=200)",     "565 ms",                   "in-process bench, slide 5"],
        ["FP8 + prefix caching e2e",         "455 ms",                   "verdict table, slide 8"],
        ["Cost per baseline request",         "$0.002303 × 0.565 = $0.001301",  "arithmetic"],
        ["Cost per combo request",            "$0.002303 × 0.455 = $0.001047",  "arithmetic"],
        ["Per 1K requests, baseline → combo", "$1.30 → $1.05  (−$0.25)",  "× 1000"],
        ["Annual savings at 1M req/day",     "$0.000254 × 1M × 365 ≈ $93K / pod / year", "arithmetic"],
    ]
    add_table(s, rows, left=0.4, top=1.95, width=12.5, height=4.6, pt=12,
              col_widths=[4.6, 4.4, 3.5])
    add_text(s,
             "Assumptions: pod owns the whole g7e.12xlarge instance; on-demand "
             "list price (reserved would be lower); chat-app-typical request "
             "shape (1024-token input, 200-token output). Annual figure scales "
             "linearly with traffic.",
             left=0.4, top=6.65, width=12.5, height=0.5,
             pt=10, italic=True, color=SUBTLE)
    add_footer(s, 22)


def appendix_flop_math(prs: Presentation) -> None:
    s = blank(prs)
    add_title(s, "Why prefill isn't quadratic in input length",
              eyebrow="APPENDIX")
    add_picture(s, FIG_DIR / "appendix_flop_math.png",
                left=2.0, top=1.95, width=9.0, height=4.85)
    add_footer(s, 23)
    set_notes(s,
        "The textbook claim that 'attention is O(N²)' is true asymptotically "
        "but misleading at chat-app scale. The full walkthrough if asked:\n\n"
        "FOUR NUMBERS THAT DESCRIBE OLMoE'S SHAPE:\n\n"
        "d_model = 2048. Each token is represented inside the model as a "
        "vector of 2048 numbers. Pick any token in your prompt — right "
        "now, inside the model, it's 2048 floats.\n\n"
        "num_heads = 16, d_head = 128. Attention doesn't do one big "
        "computation; it splits the 2048 dimensions into 16 parallel "
        "'heads' of 128 each, so 16 × 128 = 2048. Each head can learn to "
        "pay attention to a different kind of thing (one head for "
        "grammatical subjects, one for negation, one for names, etc.). "
        "Splitting into heads doesn't change the total work — it just "
        "parallelizes it across heads that learn different patterns.\n\n"
        "d_ff ≈ 1024. Each of OLMoE's experts is a small feed-forward "
        "network — two matmuls back-to-back with an activation in "
        "between. It takes the token's 2048-dimensional representation, "
        "projects it into an internal 1024-dimensional space, and "
        "projects it back out to 2048. The internal dimension d_ff is "
        "how much 'thinking room' each expert has. In OLMoE's MoE "
        "design, experts are deliberately narrow (d_ff < d_model) "
        "because you're activating 8 of them per token — they're "
        "specialists, and each one does less work than a traditional "
        "dense feed-forward network would.\n\n"
        "THE ONE FORMULA: a matrix multiplication of shape "
        "A (m × k) · B (k × n) costs about 2 · m · n · k floating-point "
        "operations. Every matmul in a transformer fits this shape. If "
        "you remember only one formula from this talk, pick this one — "
        "you can estimate any model's compute cost on a napkin.\n\n"
        "COUNT PER LAYER AT N=1024:\n\n"
        "Attention's expensive piece is the score matrix — each of N "
        "tokens takes a dot product with every other token. That's two "
        "back-to-back matmuls of shape (N × d_head) · (d_head × N) and "
        "its transpose, summed across all 16 heads. Per layer:\n"
        "  Attention FLOPs = 4 · N² · (num_heads · d_head)\n"
        "                  = 4 · N² · d_model\n"
        "                  = 4 · (1024)² · 2048\n"
        "                  ≈ 8.6 × 10⁹  =  8.6 GFLOPs\n\n"
        "MoE FFN's expensive piece is a wide matmul per active expert. "
        "Each token activates 8 experts; each expert has the "
        "2048 → 1024 → 2048 shape, which is two matmuls. Per token:\n"
        "  FFN FLOPs per token = 8 experts × 2 matmuls × 2 · d_model · d_ff\n"
        "                      = 8 × 2 × 2 · 2048 · 1024\n"
        "                      ≈ 6.7 × 10⁷\n"
        "Multiply by N tokens in the prefill:\n"
        "  FFN FLOPs per layer = (6.7 × 10⁷) · N\n"
        "                      = (6.7 × 10⁷) · 1024\n"
        "                      ≈ 6.9 × 10¹⁰  =  69 GFLOPs\n\n"
        "SIDE BY SIDE AT N=1024:\n"
        "  Attention:  8.6 GFLOPs   (grows as N²)\n"
        "  MoE FFN:    69  GFLOPs   (grows as N, with a huge coefficient)\n\n"
        "FFN does 8× more work than attention at this length. Yes, "
        "attention is asymptotically worse — its growth rate eventually "
        "wins. But look at the starting point: the FFN coefficient is "
        "~8,000 times larger than the attention coefficient. It takes a "
        "while for quadratic growth to catch up with a head start that "
        "big.\n\n"
        "CROSSOVER: attention equals FFN when N ≈ 8 · d_ff ≈ 8,000 "
        "tokens. Below that, FFN dominates. Above that, attention "
        "dominates.\n\n"
        "WHY THIS MATTERS FOR OLMoE: the model's max context is 4,096 "
        "tokens. The crossover is past what the model can process in a "
        "single request. For every request this model can serve, prefill "
        "is FFN-dominated. TTFT scales roughly linearly with input "
        "length, exactly as the measured grid (slide 5) shows.\n\n"
        "GENERALIZATION: most modern transformers have wide FFNs and "
        "modest context windows, so they live in the FFN-dominated "
        "regime during prefill. The 'attention is N²' intuition only "
        "kicks in at very long contexts (32K+) where the score matrix "
        "gets big enough to win on its coefficient. At chat-app input "
        "lengths it's all FFN.\n\n"
        "PRINCIPLE: Big-O is asymptotic. Coefficients win until you "
        "reach the asymptote. For this model at this scale, you're "
        "nowhere near the asymptote where N² would dominate."
    )


# ------------------------------------------------------------
# Build
# ------------------------------------------------------------

def build() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Spine — 15 slides
    slide_title(prs)
    slide_problem(prs)
    slide_system(prs)
    slide_method(prs)
    slide_baseline(prs)
    slide_profile(prs)
    slide_lever_board(prs)
    slide_sla_verdict_table_recommended(prs)
    slide_sla_verdict_table_all_on(prs)
    slide_sla_verdict_chart_ttft(prs)
    slide_sla_verdict_chart_tpot(prs)
    slide_sla_verdict_chart_e2e(prs)
    slide_roi(prs)
    slide_architecture_tradeoffs(prs)
    slide_recommendation(prs)
    slide_scaling_remeasure(prs)
    slide_repo(prs)
    slide_code_overview(prs)
    slide_qa(prs)

    # Appendix — backup material
    slide_business_impact(prs)   # continuous batching, throughput vs latency
    appendix_fp8_accuracy(prs)
    appendix_cost_math(prs)
    appendix_flop_math(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
